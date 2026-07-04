import csv
import re
import tempfile
from pathlib import Path
from typing import Any

import fitz
from docx import Document as DocxDocument
from docx.oxml.ns import qn
from minio import Minio
from neo4j_graphrag.experimental.components.text_splitters.fixed_size_splitter import (
    FixedSizeSplitter,
)
from neo4j_graphrag.experimental.components.types import TextChunk, TextChunks
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader

from .config import settings
from .docx_loader import load_docx_pages
from .models import ExtractRequest, SourcePage
from .visual_types import VisualCandidate


VISUAL_CAPTION = re.compile(
    r"(?iu)\b(рис(?:унок)?\.?|табл(?:ица)?\.?|figure|fig\.?|table|chart|diagram)"
)
IMAGE_CAPTION = re.compile(
    r"(?iu)\b(рис(?:унок)?\.?|figure|fig\.?)"
)

minio_client = Minio(
    settings.minio_endpoint.replace("http://", "").replace("https://", ""),
    access_key=settings.minio_access_key,
    secret_key=settings.minio_secret_key,
    secure=settings.minio_endpoint.startswith("https://"),
)


def load_source_content(
    request: ExtractRequest,
) -> tuple[list[SourcePage], list[VisualCandidate]]:
    suffix = "." + request.type.lower().lstrip(".")
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        minio_client.fget_object(settings.minio_bucket, request.storageKey, tmp.name)
        path = Path(tmp.name)

    try:
        if request.type == "pdf":
            return load_pdf_content(path)
        if request.type == "docx":
            pages = load_docx_pages(path)
            return pages, load_docx_visuals(path, pages)
        if request.type == "pptx":
            return load_pptx_content(path)
        if request.type == "xlsx":
            return load_xlsx_content(path)
        if request.type == "csv":
            return load_csv_content(path)
        return (
            [
                SourcePage(
                    page=1,
                    text=path.read_text(encoding="utf-8", errors="ignore"),
                    section="Документ",
                )
            ],
            [],
        )
    finally:
        path.unlink(missing_ok=True)


def load_source_pages(request: ExtractRequest) -> list[SourcePage]:
    pages, _ = load_source_content(request)
    return pages


def load_pdf_content(
    path: Path,
) -> tuple[list[SourcePage], list[VisualCandidate]]:
    pages: list[SourcePage] = []
    candidates: list[VisualCandidate] = []
    try:
        document = fitz.open(str(path))
        for index, page in enumerate(document, start=1):
            text = page.get_text("text") or ""
            pages.append(
                SourcePage(
                    page=index,
                    text=text,
                    section=f"Страница {index}",
                )
            )
            candidates.extend(extract_pdf_tables(page, index))
            image_count = len(page.get_images(full=True))
            drawing_count = len(page.get_drawings())
            if (
                image_count > 0
                or drawing_count >= 6
                or VISUAL_CAPTION.search(text)
            ):
                pixmap = page.get_pixmap(
                    matrix=fitz.Matrix(1.35, 1.35),
                    alpha=False,
                )
                candidates.append(
                    VisualCandidate(
                        kind="image",
                        page=index,
                        section=f"Страница {index}",
                        title=first_visual_caption(text)
                        or f"Визуальное содержимое страницы {index}",
                        text=text[:1200],
                        image=pixmap.tobytes("jpeg"),
                        content_type="image/jpeg",
                    )
                )
        document.close()
        return pages, candidates
    except Exception:
        # Preserve text ingestion for encrypted and unusual PDFs.
        return (
            [
                SourcePage(
                    page=index,
                    text=page.extract_text() or "",
                    section=f"Страница {index}",
                )
                for index, page in enumerate(PdfReader(str(path)).pages, start=1)
            ],
            [],
        )


def extract_pdf_tables(page: fitz.Page, page_number: int) -> list[VisualCandidate]:
    try:
        tables = page.find_tables().tables
    except Exception:
        return []
    result: list[VisualCandidate] = []
    for index, table in enumerate(tables, start=1):
        markdown = rows_to_markdown(table.extract())
        if markdown:
            result.append(
                VisualCandidate(
                    kind="table",
                    page=page_number,
                    section=f"Страница {page_number}",
                    title=f"Таблица {index} на странице {page_number}",
                    text=markdown,
                )
            )
    return result


def load_docx_visuals(
    path: Path,
    pages: list[SourcePage],
) -> list[VisualCandidate]:
    document = DocxDocument(str(path))
    candidates: list[VisualCandidate] = []
    for index, table in enumerate(document.tables, start=1):
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        markdown = rows_to_markdown(rows)
        if markdown:
            candidates.append(
                VisualCandidate(
                    kind="table",
                    page=find_source_page(
                        pages,
                        rows[0][0] if rows and rows[0] else "",
                    ),
                    section="Таблица",
                    title=f"Таблица {index}",
                    text=markdown,
                )
            )

    seen_relationships: set[str] = set()
    image_captions = [
        (page.page, line.strip())
        for page in pages
        for line in page.text.splitlines()
        if IMAGE_CAPTION.search(line)
    ]
    for node in document.element.body.xpath(".//*[local-name()='blip']"):
        relationship_id = node.get(qn("r:embed"))
        if not relationship_id or relationship_id in seen_relationships:
            continue
        seen_relationships.add(relationship_id)
        part = document.part.related_parts.get(relationship_id)
        blob = getattr(part, "blob", None)
        content_type = getattr(part, "content_type", None)
        if not blob or not str(content_type or "").startswith("image/"):
            continue
        caption_index = len(seen_relationships) - 1
        caption = (
            image_captions[caption_index]
            if caption_index < len(image_captions)
            else None
        )
        candidates.append(
            VisualCandidate(
                kind="image",
                page=caption[0] if caption else 1,
                section="Изображение",
                title=(
                    caption[1]
                    if caption
                    else f"Изображение {len(seen_relationships)}"
                ),
                image=blob,
                content_type=content_type,
            )
        )
    return candidates


def load_pptx_content(
    path: Path,
) -> tuple[list[SourcePage], list[VisualCandidate]]:
    presentation = Presentation(str(path))
    pages: list[SourcePage] = []
    candidates: list[VisualCandidate] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        text_items: list[str] = []
        diagram_shape_count = 0
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = str(getattr(shape, "text", "") or "").strip()
                if text:
                    text_items.append(text)
            if getattr(shape, "has_table", False):
                rows = [
                    [cell.text for cell in row.cells]
                    for row in shape.table.rows
                ]
                markdown = rows_to_markdown(rows)
                if markdown:
                    candidates.append(
                        VisualCandidate(
                            kind="table",
                            page=slide_number,
                            section=f"Слайд {slide_number}",
                            title=f"Таблица на слайде {slide_number}",
                            text=markdown,
                        )
                    )
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                candidates.append(
                    VisualCandidate(
                        kind="image",
                        page=slide_number,
                        section=f"Слайд {slide_number}",
                        title=str(
                            getattr(shape, "name", "") or "Изображение"
                        ),
                        image=shape.image.blob,
                        content_type=shape.image.content_type,
                    )
                )
            if getattr(shape, "has_chart", False):
                candidates.append(
                    VisualCandidate(
                        kind="chart",
                        page=slide_number,
                        section=f"Слайд {slide_number}",
                        title=f"График на слайде {slide_number}",
                        text=chart_to_markdown(shape.chart),
                    )
                )
            if shape.shape_type in {
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                MSO_SHAPE_TYPE.GROUP,
                MSO_SHAPE_TYPE.LINE,
            }:
                diagram_shape_count += 1
        if diagram_shape_count >= 4 and len(text_items) >= 3:
            candidates.append(
                VisualCandidate(
                    kind="diagram",
                    page=slide_number,
                    section=f"Слайд {slide_number}",
                    title=f"Схема на слайде {slide_number}",
                    text=(
                        "Элементы схемы:\n- "
                        + "\n- ".join(text_items)
                    ),
                )
            )
        pages.append(
            SourcePage(
                page=slide_number,
                text="\n".join(text_items),
                section=f"Слайд {slide_number}",
            )
        )
    return pages, candidates


def load_xlsx_content(
    path: Path,
) -> tuple[list[SourcePage], list[VisualCandidate]]:
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    pages: list[SourcePage] = []
    candidates: list[VisualCandidate] = []
    for index, sheet in enumerate(workbook.worksheets, start=1):
        rows = [list(row) for row in sheet.iter_rows(values_only=True)]
        markdown = rows_to_markdown(rows)
        pages.append(SourcePage(page=index, text=markdown, section=sheet.title))
        if markdown:
            candidates.append(
                VisualCandidate(
                    kind="table",
                    page=index,
                    section=sheet.title,
                    title=f"Лист «{sheet.title}»",
                    text=markdown,
                )
            )
    workbook.close()
    return pages, candidates


def load_csv_content(
    path: Path,
) -> tuple[list[SourcePage], list[VisualCandidate]]:
    with path.open("r", encoding="utf-8", errors="ignore") as file:
        rows = list(csv.reader(file))
    markdown = rows_to_markdown(rows)
    page = SourcePage(page=1, text=markdown, section="CSV")
    candidate = VisualCandidate(
        kind="table",
        page=1,
        section="CSV",
        title="Таблица CSV",
        text=markdown,
    )
    return [page], [candidate] if markdown else []


async def split_source_pages(
    document_id: str,
    pages: list[SourcePage],
) -> TextChunks:
    splitter = FixedSizeSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
        approximate=False,
    )
    chunks: list[TextChunk] = []
    global_index = 0
    for page in pages:
        if not page.text.strip():
            continue
        page_chunks = await splitter.run(text=page.text)
        for page_chunk in page_chunks.chunks:
            metadata: dict[str, Any] = {
                "page": page.page,
                "section": page.section or "",
                "documentId": document_id,
            }
            if page.visualId:
                metadata["visualId"] = page.visualId
                metadata["visualType"] = page.visualType or "image"
            chunks.append(
                TextChunk(
                    uid=f"{document_id}-chunk-{global_index}",
                    index=global_index,
                    text=page_chunk.text,
                    metadata=metadata,
                )
            )
            global_index += 1
    return TextChunks(chunks=chunks)


def rows_to_markdown(rows: list[list[Any]]) -> str:
    normalized = [
        [normalize_cell(cell) for cell in row]
        for row in rows
        if any(normalize_cell(cell) for cell in row)
    ]
    if not normalized:
        return ""
    width = max(len(row) for row in normalized)
    padded = [row + [""] * (width - len(row)) for row in normalized]
    lines = [
        "| " + " | ".join(padded[0]) + " |",
        "| " + " | ".join("---" for _ in range(width)) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in padded[1:])
    return "\n".join(lines)


def normalize_cell(value: Any) -> str:
    if value is None:
        return ""
    return re.sub(r"\s+", " ", str(value)).strip().replace("|", "\\|")


def first_visual_caption(text: str) -> str | None:
    for line in text.splitlines():
        normalized = line.strip()
        if normalized and VISUAL_CAPTION.search(normalized):
            return normalized[:240]
    return None


def find_source_page(pages: list[SourcePage], marker: str) -> int:
    normalized = re.sub(r"\s+", " ", marker).strip()
    if normalized:
        for page in pages:
            if normalized[:80] in re.sub(r"\s+", " ", page.text):
                return page.page
    return 1


def chart_to_markdown(chart: Any) -> str:
    try:
        series_items = list(chart.series)
        categories = list(chart.plots[0].categories)
        if categories:
            rows: list[list[Any]] = [
                [
                    "Категория",
                    *[
                        str(getattr(series, "name", "") or f"Серия {index}")
                        for index, series in enumerate(series_items, start=1)
                    ],
                ]
            ]
            for index, category in enumerate(categories):
                rows.append(
                    [
                        str(getattr(category, "label", category)),
                        *[
                            (
                                series.values[index]
                                if index < len(series.values)
                                else ""
                            )
                            for series in series_items
                        ],
                    ]
                )
            return rows_to_markdown(rows)
        rows = [["Серия", "Значения"]]
        for series in series_items:
            rows.append([
                str(getattr(series, "name", "") or "Серия"),
                ", ".join(str(value) for value in series.values),
            ])
    except Exception:
        return "График без доступных структурированных данных."
    return rows_to_markdown(rows)
